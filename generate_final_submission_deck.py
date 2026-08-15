import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_final_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Theme colors
    BG_DARK = RGBColor(15, 23, 42)       # #0F172A Deep Slate / Dark Navy
    CARD_BG = RGBColor(30, 41, 59)       # #1E293B Card Background
    CARD_BORDER = RGBColor(51, 65, 85)   # #334155 Card Border
    TEXT_MAIN = RGBColor(248, 250, 252)  # #F8FAFC Off-White Text
    TEXT_MUTED = RGBColor(148, 163, 184) # #94A3B8 Muted Slate Text
    ACCENT_CYAN = RGBColor(6, 182, 212)  # #06B6D4 Electric Cyan
    ACCENT_GREEN = RGBColor(16, 185, 129)# #10B981 Emerald Green
    ACCENT_BLUE = RGBColor(59, 130, 246) # #3B82F6 Bright Blue
    ACCENT_AMBER = RGBColor(245, 158, 11)# #F59E0B Warm Amber

    def apply_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()

    def add_header(slide, title_text, category_text="PROJECT SUBMISSION DECK"):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_CYAN
        p_cat.font.name = "Arial"

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_MAIN
        p_title.font.name = "Arial"

    def add_card(slide, left, top, width, height, title, items, border_color=CARD_BORDER, title_color=ACCENT_CYAN):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.20)
        tf.margin_bottom = Inches(0.20)

        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(17)
        p0.font.bold = True
        p0.font.color.rgb = title_color
        p0.font.name = "Arial"
        p0.space_after = Pt(10)

        for item in items:
            p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_MAIN
            p.font.name = "Arial"
            p.space_after = Pt(6)

    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Submission Version)
    # ----------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    apply_background(s1)

    tb1 = s1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "⚡ ELECTRICITY THEFT & ANOMALY DETECTION SYSTEM"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(12)

    p2 = tf1.add_paragraph()
    p2.text = "An End-to-End Decision Support System Powered by Dual-Signal ML & Explainable AI"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_MAIN
    p2.space_after = Pt(20)

    p3 = tf1.add_paragraph()
    p3.text = "Final Project Submission Deck | Beginner-Friendly Guide & Full-Stack Documentation"
    p3.font.size = Pt(14)
    p3.font.bold = True
    p3.font.color.rgb = ACCENT_GREEN
    p3.space_after = Pt(12)

    p4 = tf1.add_paragraph()
    p4.text = "Includes Core Problem Explanation, Dual-Signal ML Engine, Dedicated Backend Architecture, Dedicated Frontend UI, Datasets, and Real-World Guidelines."
    p4.font.size = Pt(13)
    p4.font.color.rgb = TEXT_MUTED

    s1.notes_slide.notes_text_frame.text = (
        "Submission Documentation: Welcome to the final submission deck for the Electricity Theft Detection System. "
        "This presentation is structured so that any reviewer—technical or non-technical—can easily understand the real-world problem, "
        "the machine learning innovations, and the exact backend & frontend software implementations."
    )

    # ----------------------------------------------------
    # SLIDE 2: Understanding Electricity Theft (Beginner-Friendly Problem)
    # ----------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    apply_background(s2)
    add_header(s2, "1. Understanding Electricity Theft: What & Why?")

    add_card(s2, 0.8, 1.8, 3.6, 5.0, "What is Electricity Theft?", [
        "Non-Technical Losses (NTL): Power consumed but never billed.",
        "Physical Meter Tampering: Slowing down mechanical/digital meters.",
        "Meter Bypass: Connecting lines directly to bypass the meter entirely.",
        "Cyber Alteration: Hacking smart meter communication feeds."
    ], title_color=ACCENT_AMBER)

    add_card(s2, 4.8, 1.8, 3.6, 5.0, "Why is it a Massive Problem?", [
        "Economic Loss: $96B+ lost globally every year by utility companies.",
        "Grid Unreliability: Causes unexpected transformer overloads and power outages.",
        "Higher Bills: Paying honest consumers subsidize stolen electricity.",
        "Safety Risk: Illegal wiring creates severe electrical fire hazards."
    ], title_color=ACCENT_AMBER)

    add_card(s2, 8.8, 1.8, 3.6, 5.0, "Why is it Hard to Catch?", [
        "Data Scale: Millions of households streaming daily electricity readings.",
        "Physical Inspections: Checking every meter manually is physically impossible.",
        "Evolving Tricks: Thieves constantly change their tampering methods.",
        "Need for High Precision: False allegations damage utility customer trust."
    ], title_color=ACCENT_AMBER)

    s2.notes_slide.notes_text_frame.text = (
        "Slide 2 Context: Electricity theft is a huge global issue where power is consumed illegally without billing. "
        "It leads to massive financial losses ($96B+/year) and grid instability. Because power companies monitor millions of meters daily, "
        "manual physical inspections are too slow and expensive."
    )

    # ----------------------------------------------------
    # SLIDE 3: Why Traditional Detection Methods Fail
    # ----------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    apply_background(s3)
    add_header(s3, "2. Why Traditional Detection Methods Fail")

    add_card(s3, 0.8, 1.8, 3.6, 5.0, "1. Historical Label Bias", [
        "Standard models train ONLY on past confirmed theft cases.",
        "Historical theft labels are sparse and incomplete.",
        "Models learn ONLY what inspectors caught in the past.",
        "Completely blind to novel ('zero-day') tampering tricks."
    ])

    add_card(s3, 4.8, 1.8, 3.6, 5.0, "2. Black-Box Score Trap", [
        "Standard ML outputs a single raw score (e.g. '89% suspicious').",
        "Provides ZERO explanation of why the score is high.",
        "Inspectors hesitate to visit sites without evidence.",
        "Wastes field budget on false alarms."
    ])

    add_card(s3, 8.8, 1.8, 3.6, 5.0, "3. Weather & Seasonal Confusion", [
        "A sudden drop in usage might be cold weather, not theft.",
        "If an entire neighborhood drops consumption, it's seasonal.",
        "Single-meter models fail to compare against neighbors.",
        "Triggers frequent false accusations during mild weather."
    ])

    s3.notes_slide.notes_text_frame.text = (
        "Slide 3 Context: Conventional machine learning tutorials simply train one classifier on past theft labels. "
        "This fails in practice because: (1) past labels inherit human inspector bias, (2) raw scores lack explainable proof, "
        "and (3) simple models mistake seasonal cold weather drops for meter tampering."
    )

    # ----------------------------------------------------
    # SLIDE 4: The Solution Overview (Dual-Signal & Explainable AI)
    # ----------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    apply_background(s4)
    add_header(s4, "3. Our Solution: Dual-Signal ML + SHAP Explainability")

    add_card(s4, 0.8, 1.8, 5.6, 2.4, "Signal 1: Supervised XGBoost", [
        "Catches KNOWN historical theft patterns.",
        "Uses class-weighted gradient boosting to handle sparse theft cases.",
        "Calculates precise supervised probability (P_supervised)."
    ], title_color=ACCENT_BLUE)

    add_card(s4, 6.9, 1.8, 5.6, 2.4, "Signal 2: Unsupervised Isolation Forest", [
        "Catches NOVEL / ZERO-DAY tampering anomalies.",
        "Requires NO historical theft labels to detect strange usage.",
        "Calculates statistical outlier score (Score_anomaly)."
    ], title_color=ACCENT_GREEN)

    add_card(s4, 0.8, 4.4, 5.6, 2.4, "Signal 3: SHAP Explainability Engine", [
        "Computes game-theoretic feature attributions.",
        "Translates mathematical weights into plain-English reasons.",
        "Tells inspectors EXACTLY why an account was flagged."
    ], title_color=ACCENT_CYAN)

    add_card(s4, 6.9, 4.4, 5.6, 2.4, "Signal 4: Neighborhood Peer Divergence", [
        "Compares consumer usage against neighbors on the same transformer.",
        "Filters out weather/seasonal drops affecting the whole feeder.",
        "Pinpoints true divergence vs. normal neighborhood behavior."
    ], title_color=ACCENT_AMBER)

    s4.notes_slide.notes_text_frame.text = (
        "Slide 4 Context: Our solution combines four intelligent signals: (1) Supervised XGBoost for known theft, "
        "(2) Unsupervised Isolation Forest for unseen zero-day anomalies, (3) SHAP explainability for plain-English proof, "
        "and (4) Peer divergence to filter out weather effects."
    )

    # ----------------------------------------------------
    # SLIDE 5: DEDICATED SLIDE - BACKEND ARCHITECTURE
    # ----------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    apply_background(s5)
    add_header(s5, "4. Deep Dive: Backend Architecture & Implementation")

    add_card(s5, 0.8, 1.8, 5.6, 5.0, "⚙️ Backend Core Modules (src/)", [
        "1. Feature Engine (src/features.py): Computes 15 domain features per consumer series (distribution, zero streaks, z-scores).",
        "2. ML Ensemble (src/models.py): Implements TheftDetectionEnsemble combining XGBoost + Isolation Forest.",
        "3. SHAP Engine (src/explain.py): SHAP TreeExplainer generating positive & negative feature attributions.",
        "4. Data Pipeline (src/train.py & generate_data.py): Automated synthetic generator & model trainer."
    ], title_color=ACCENT_BLUE)

    add_card(s5, 6.9, 1.8, 5.6, 5.0, "🔌 FastAPI REST Endpoints (src/api.py)", [
        "• POST /scan : Batch population scanner. Ranks all grid consumers by risk score and returns top N suspects.",
        "• POST /score : Ad-hoc single consumer scorer. Takes raw daily kWh array and computes real-time risk & reasons.",
        "• High Performance: Asynchronous execution, strict Pydantic validation schemas (`schemas.py`), production ready."
    ], title_color=ACCENT_CYAN)

    s5.notes_slide.notes_text_frame.text = (
        "Slide 5 Backend Focus: This slide is 100% dedicated to the Backend. "
        "The backend is built in Python using FastAPI. It consists of modular Python scripts for feature engineering, "
        "dual-signal ensemble scoring, SHAP reason generation, and high-performance REST API endpoints."
    )

    # ----------------------------------------------------
    # SLIDE 6: DEDICATED SLIDE - FRONTEND OPERATOR UI
    # ----------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    apply_background(s6)
    add_header(s6, "5. Deep Dive: Frontend UX & Operator Interface")

    add_card(s6, 0.8, 1.8, 5.6, 5.0, "🖥️ Streamlit Web Application (dashboard.py)", [
        "1. Real-Time API Integration: Connects via HTTP to FastAPI backend (`THEFT_API_URL`).",
        "2. Sidebar Scan Control: Interactive slider (top 5 to 100 suspects) & primary trigger button.",
        "3. Top Executive KPI Cards: Displays Monitored Consumers, Flagged High-Risk Cases & Active Risk Threshold.",
        "4. Ranked Suspect Table: Interactive grid showing Consumer ID, Transformer Zone, Overall Risk %, ML Prob % & Anomaly %."
    ], title_color=ACCENT_GREEN)

    add_card(s6, 6.9, 1.8, 5.6, 5.0, "📊 Visual Analytics & Audit Inspector", [
        "1. Transformer Cluster Chart: Plotly bar chart displaying suspect counts grouped by transformer zone.",
        "2. Automated Audit Log Inspector: Consumer dropdown allowing instant deep-dive into flagged accounts.",
        "3. Plain-English Reason Render: Renders red alert callouts (🛑) detailing exact failure modes.",
        "4. Inspector Decision Support: Gives field teams visual proof before dispatching on-site audits."
    ], title_color=ACCENT_AMBER)

    s6.notes_slide.notes_text_frame.text = (
        "Slide 6 Frontend Focus: This slide is 100% dedicated to the Frontend. "
        "The frontend is built using Streamlit and Plotly. It provides a clean, responsive operator dashboard "
        "where utility managers can trigger live grid threat scans, visualize suspect clusters by transformer, and inspect automated audit logs."
    )

    # ----------------------------------------------------
    # SLIDE 7: Domain Feature Engineering (15 Smart Indicators)
    # ----------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    apply_background(s7)
    add_header(s7, "6. Feature Engineering & Domain Intelligence (15 Features)")

    add_card(s7, 0.8, 1.8, 3.6, 5.0, "Consumption Dynamics", [
        "Mean, Std, Skewness, Kurtosis",
        "Coefficient of Variation (CV)",
        "Normalized Trend Slope",
        "Largest Single-Day Drop %",
        "Count of Sudden Drops"
    ])

    add_card(s7, 4.8, 1.8, 3.6, 5.0, "Bypass & Periodicity", [
        "Zero-Consumption Ratio",
        "Longest Zero Streak (meter bypass signature)",
        "Weekly Autocorrelation",
        "Weekday vs Weekend Ratio (rhythm disruption)"
    ])

    add_card(s7, 8.8, 1.8, 3.6, 5.0, "Peer Divergence", [
        "Transformer-Group Z-Score",
        "Feeder Correlation Score",
        "Controls for weather & season by comparing against immediate transformer neighbors."
    ])

    s7.notes_slide.notes_text_frame.text = (
        "Slide 7 Feature Focus: We extract 15 domain features from raw daily meter readings. "
        "These include distribution stats, meter bypass streaks, weekly autocorrelation drops, and crucial transformer peer divergence scores."
    )

    # ----------------------------------------------------
    # SLIDE 8: Risk Scoring & Disambiguation Matrix
    # ----------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    apply_background(s8)
    add_header(s8, "7. Risk Scoring Formula & Disambiguation Matrix")

    add_card(s8, 0.8, 1.8, 11.7, 1.5, "Composite Risk Score Formula", [
        "Risk Score = 0.65 * Supervised_Prob + 0.35 * Anomaly_Score",
        "Supervised and Unsupervised components are reported separately to enable precise decision support."
    ])

    add_card(s8, 0.8, 3.6, 5.6, 3.2, "High Supervised + High Anomaly", [
        "CLASSIFICATION: High-Confidence Known Theft",
        "ACTION: High-priority dispatch for physical site audit.",
        "SIGNATURE: Known tampering pattern with strong statistical outlier metrics."
    ])

    add_card(s8, 6.9, 3.6, 5.6, 3.2, "Low Supervised + High Anomaly", [
        "CLASSIFICATION: Zero-Day Anomaly / Novel Pattern",
        "ACTION: Flagged for Human Utility Review (Not automatic theft).",
        "SIGNATURE: Unusual usage (vacant home or new tampering method)."
    ])

    s8.notes_slide.notes_text_frame.text = (
        "Slide 8 Scoring Matrix: We blend supervised probability (65%) and anomaly score (35%). "
        "Reporting both components separately prevents false accusations—an account with high anomaly but low supervised probability "
        "is flagged as 'needs human review' rather than a definitive theft call."
    )

    # ----------------------------------------------------
    # SLIDE 9: Datasets & Industry Benchmarks
    # ----------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    apply_background(s9)
    add_header(s9, "8. Datasets: Synthetic Quickstart & Real Benchmarks")

    add_card(s9, 0.8, 1.8, 5.6, 5.0, "Synthetic Data Generator", [
        "Built-in generator for instant zero-dependency setup.",
        "Simulates normal load profiles, drop-offs, and bypasses.",
        "Used for automated CI/CD unit testing & quick demo verification.",
        "Reaches high ROC-AUC on clean benchmark scenarios."
    ])

    add_card(s9, 6.9, 1.8, 5.6, 5.0, "Real-World Public Benchmarks", [
        "SGCC Benchmark (State Grid Corp of China):",
        "  • 42,372 consumers, 3,615 confirmed theft cases (IEEE TII standard).",
        "  • Plug-and-play format via simple melt transformation.",
        "PRECON & Irish CER Datasets:",
        "  • High-frequency smart meter data for 30-minute interval analysis."
    ])

    s9.notes_slide.notes_text_frame.text = (
        "Slide 9 Datasets: The system includes a synthetic generator for zero-dependency testing, "
        "and natively supports real-world industry benchmarks like the 42,000-consumer SGCC dataset."
    )

    # ----------------------------------------------------
    # SLIDE 10: Honest Real-World Caveats
    # ----------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    apply_background(s10)
    add_header(s10, "9. Honest Real-World Caveats & Operational Guidelines")

    add_card(s10, 0.8, 1.8, 5.6, 5.0, "Technical & Data Realities", [
        "Synthetic vs. Real Noise: Synthetic data is cleanly separable (~1.0 AUC); real-world grid data is significantly messier.",
        "Topology Reliance: Peer comparisons rely on accurate transformer mapping. Bad feeder data affects z-scores.",
        "Vacant Property Outliers: Anomaly detection flags vacant/relocated homes—hence the need for human review."
    ])

    add_card(s10, 6.9, 1.8, 5.6, 5.0, "Operational & Policy Guidelines", [
        "Decision Support Tool: The system assists human inspectors; it is NOT an automated legal/billing penalty engine.",
        "Feedback Loop: Field inspection results feed back into retraining the supervised XGBoost model.",
        "Privacy & Governance: Operates on anonymized consumption series."
    ])

    s10.notes_slide.notes_text_frame.text = (
        "Slide 10 Caveats: Transparency is vital. Real grid data is messier than synthetic generators, "
        "peer comparison requires correct feeder topology, and our tool is designed for human decision support, not automatic billing fines."
    )

    # ----------------------------------------------------
    # SLIDE 11: Future Roadmap
    # ----------------------------------------------------
    s11 = prs.slides.add_slide(blank_layout)
    apply_background(s11)
    add_header(s11, "10. Future Roadmap & Technical Enhancements")

    add_card(s11, 0.8, 1.8, 5.6, 5.0, "Algorithmic & Model Extensions", [
        "LightGBM / CatBoost: Swap XGBoost for faster training on multi-million consumer grids.",
        "Deep Learning Waveform Autoencoders: Add 1D-CNN / LSTM autoencoders on raw 30-min series as a 3rd signal.",
        "Advanced Imbalance Handling: Integrate SMOTE (imbalanced-learn) for grid datasets with <1% theft ratio."
    ], title_color=ACCENT_BLUE)

    add_card(s11, 6.9, 1.8, 5.6, 5.0, "Platform & Persistence Enhancements", [
        "Risk Trajectory Persistence: SQLite scan history tracking rising risk score trends over successive billing cycles.",
        "GIS & Map Integration: Overlay transformer suspect clusters onto spatial GIS map layers.",
        "Automated Inspector Mobile App: Mobile API integration for field inspection teams."
    ], title_color=ACCENT_GREEN)

    s11.notes_slide.notes_text_frame.text = (
        "Slide 11 Roadmap: Outlines future scalability path including LightGBM/CatBoost swaps, "
        "deep learning waveform autoencoders, SQLite risk trajectory persistence, and mobile GIS field integration."
    )

    # ----------------------------------------------------
    # SLIDE 12: Executive Submission Summary
    # ----------------------------------------------------
    s12 = prs.slides.add_slide(blank_layout)
    apply_background(s12)
    add_header(s12, "11. Final Submission Summary & Key Achievements")

    add_card(s12, 0.8, 1.8, 5.6, 5.0, "🎯 Problem & Algorithmic Impact", [
        "Solves Real Problem: Reduces $96B+ NTL losses for power utilities.",
        "Dual-Signal ML: XGBoost (known theft) + Isolation Forest (zero-day anomalies).",
        "Peer Divergence: Eliminates weather false alarms via transformer z-scores.",
        "Explainable AI: SHAP reasons restore human inspector trust."
    ], title_color=ACCENT_CYAN)

    add_card(s12, 6.9, 1.8, 5.6, 5.0, "🚀 Full-Stack Software Excellence", [
        "Modular Backend: FastAPI REST service with /scan and /score endpoints.",
        "Interactive Frontend: Streamlit dashboard with Plotly geographic charts.",
        "Production Ready: Comprehensive feature extraction, schemas, unit tests.",
        "Dataset Agnostic: Supports synthetic testing & real SGCC smart meter data."
    ], title_color=ACCENT_GREEN)

    s12.notes_slide.notes_text_frame.text = (
        "Slide 12 Final Summary: In summary, our project delivers an end-to-end, production-ready Electricity Theft Detection system. "
        "It combines dual-signal ML innovation, SHAP explainability, a robust FastAPI backend, and an intuitive Streamlit frontend interface."
    )

    output_path = "Electricity_Theft_Detection_Final_Submission.pptx"
    prs.save(output_path)
    print(f"Final Submission Presentation successfully created at {output_path}")

if __name__ == "__main__":
    create_final_deck()
