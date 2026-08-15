import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Theme colors
    BG_DARK = RGBColor(15, 23, 42)       # #0F172A Deep Navy/Slate
    CARD_BG = RGBColor(30, 41, 59)       # #1E293B Card Background
    CARD_BORDER = RGBColor(51, 65, 85)   # #334155 Card Border
    TEXT_MAIN = RGBColor(248, 250, 252)  # #F8FAFC Off white
    TEXT_MUTED = RGBColor(148, 163, 184) # #94A3B8 Muted gray
    ACCENT_CYAN = RGBColor(6, 182, 212)  # #06B6D4 Electric Cyan
    ACCENT_GREEN = RGBColor(16, 185, 129)# #10B981 Emerald
    ACCENT_AMBER = RGBColor(245, 158, 11)# #F59E0B Amber
    ACCENT_BLUE = RGBColor(59, 130, 246) # #3B82F6 Bright Blue

    def apply_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()

    def add_header(slide, title_text, category_text="HACKATHON PITCH DECK"):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_CYAN
        p_cat.font.name = "Arial"

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_MAIN
        p_title.font.name = "Arial"

    def add_card(slide, left, top, width, height, title, items, border_color=CARD_BORDER, title_color=ACCENT_CYAN):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.25)
        tf.margin_bottom = Inches(0.25)

        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(18)
        p0.font.bold = True
        p0.font.color.rgb = title_color
        p0.font.name = "Arial"
        p0.space_after = Pt(12)

        for item in items:
            p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_MAIN
            p.font.name = "Arial"
            p.space_after = Pt(8)

    # ----------------------------------------------------
    # SLIDE 1: Title Slide
    # ----------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    apply_background(s1)

    tb1 = s1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "⚡ ELECTRICITY THEFT DETECTION"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(14)

    p2 = tf1.add_paragraph()
    p2.text = "Dual-Signal Machine Learning & Explainable AI (XAI) for Smart Meter Data"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_MAIN
    p2.space_after = Pt(24)

    p3 = tf1.add_paragraph()
    p3.text = "Full Stack Architecture: Modular FastAPI ML Backend + Interactive Streamlit Operator Dashboard"
    p3.font.size = Pt(16)
    p3.font.color.rgb = TEXT_MUTED

    s1.notes_slide.notes_text_frame.text = (
        "Presenter Script: Welcome judges! Today we present an end-to-end Electricity Theft Detection platform. "
        "We specifically separate our technical contributions into a high-performance Python FastAPI ML backend "
        "and a real-time Streamlit operator interface."
    )

    # ----------------------------------------------------
    # SLIDE 2: Problem Statement & Motivation
    # ----------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    apply_background(s2)
    add_header(s2, "The Problem: Flaws in Traditional Theft Detection")

    add_card(s2, 0.8, 1.8, 5.6, 5.0, "1. Sparse & Biased Ground Truth", [
        "Historical theft labels are severely incomplete.",
        "Models trained ONLY on past labels only learn what inspectors already knew.",
        "Novel meter bypasses ('zero-day theft') go completely undetected.",
        "Overfits to known fraud signatures while missing evolving tampering tricks."
    ])

    add_card(s2, 6.9, 1.8, 5.6, 5.0, "2. Black-Box Risk Scores", [
        "Utilities receive single raw probability numbers with zero context.",
        "Field inspections are costly and labor-intensive to dispatch.",
        "Without human-readable justifications, inspector trust rapidly degrades.",
        "High false-alarm rates waste field staff time and grid operational budget."
    ])

    s2.notes_slide.notes_text_frame.text = (
        "Presenter Script: Why can't we just fit XGBoost on historical labels and call it a day? "
        "First, confirmed theft labels are sparse and biased—they only represent what inspectors caught in the past. "
        "Second, utility operators need clear reasons, not black-box scores, before dispatching field crews."
    )

    # ----------------------------------------------------
    # SLIDE 3: The Hybrid Solution Overview
    # ----------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    apply_background(s3)
    add_header(s3, "The Solution: Dual-Signal ML + Explainability Engine")

    add_card(s3, 0.8, 1.8, 3.6, 5.0, "Supervised Signal", [
        "XGBoost (Class-Weighted)",
        "Detects patterns matching historical confirmed theft cases.",
        "Provides precise probability score based on domain features."
    ])

    add_card(s3, 4.8, 1.8, 3.6, 5.0, "Unsupervised Signal", [
        "Isolation Forest",
        "Flags novel, zero-day statistical anomalies.",
        "Works even without confirmed ground-truth labels."
    ])

    add_card(s3, 8.8, 1.8, 3.6, 5.0, "Explainability (XAI)", [
        "SHAP TreeExplainer",
        "Generates per-account feature attribution summaries.",
        "Tells inspectors exactly WHY an account was flagged."
    ])

    s3.notes_slide.notes_text_frame.text = (
        "Presenter Script: Our framework solves both problems by pairing a supervised XGBoost classifier with an "
        "unsupervised Isolation Forest. Plus, every prediction is powered by SHAP TreeExplainer to break down the exact drivers."
    )

    # ----------------------------------------------------
    # SLIDE 4: SEPARATE BACKEND VS FRONTEND ARCHITECTURE
    # ----------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    apply_background(s4)
    add_header(s4, "System Architecture: Backend vs. Frontend Split")

    add_card(s4, 0.8, 1.8, 5.6, 5.0, "⚙️ Backend Engine (FastAPI + ML)", [
        "Feature Engineering Engine (src/features.py): Computes 15 domain features per consumer series.",
        "Ensemble Engine (src/models.py): Combines XGBoost + Isolation Forest risk components.",
        "SHAP Reason Generator (src/explain.py): Computes exact feature attributions.",
        "REST API Service (src/api.py): Serves POST /scan (population rank) & POST /score (ad-hoc series)."
    ], title_color=ACCENT_BLUE)

    add_card(s4, 6.9, 1.8, 5.6, 5.0, "🖥️ Frontend UI (Streamlit Dashboard)", [
        "Interactive Scan Controls: Dynamic slider (top 5-100 suspects) & asynchronous trigger.",
        "Summary KPI Metrics: Displays Monitored Consumers, Flagged Cases & Risk Threshold.",
        "Ranked Suspects Table: Clean grid sorted by overall risk, ML prob & anomaly score.",
        "Geographic Cluster Chart: Plotly bar chart displaying suspect counts by transformer feeder.",
        "Audit Log Inspector: Renders automated SHAP reasons per selected consumer."
    ], title_color=ACCENT_GREEN)

    s4.notes_slide.notes_text_frame.text = (
        "Presenter Script: Notice how cleanly our solution separates responsibilities: The backend handles feature extraction, "
        "ensemble scoring, SHAP computations, and REST endpoints via FastAPI. The frontend handles interactive scan parameterization, "
        "transformer cluster visualization, and human-in-the-loop audit log inspection via Streamlit."
    )

    # ----------------------------------------------------
    # SLIDE 5: Domain Feature Engineering
    # ----------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    apply_background(s5)
    add_header(s5, "Domain-Engineered Feature Suite (15 Indicators)")

    add_card(s5, 0.8, 1.8, 3.6, 5.0, "Consumption Dynamics", [
        "Mean, Std, Skewness, Kurtosis",
        "Coefficient of Variation (CV)",
        "Normalized Trend Slope",
        "Largest Single-Day Drop %",
        "Count of Sudden Drops"
    ])

    add_card(s5, 4.8, 1.8, 3.6, 5.0, "Bypass & Periodicity", [
        "Zero-Consumption Ratio",
        "Longest Zero Streak (bypass signature)",
        "Weekly Autocorrelation",
        "Weekday vs Weekend Ratio (rhythm disruption)"
    ])

    add_card(s5, 8.8, 1.8, 3.6, 5.0, "Peer Divergence", [
        "Transformer-Group Z-Score",
        "Feeder Correlation Score",
        "Controls for weather & season by comparing against immediate transformer neighbors."
    ])

    s5.notes_slide.notes_text_frame.text = (
        "Presenter Script: Raw usage series alone isn't enough. We engineer 15 domain features spanning distribution, "
        "tamper bypass streaks, weekly periodicity disruption, and crucial transformer peer divergence."
    )

    # ----------------------------------------------------
    # SLIDE 6: Dual-Signal Scoring & Disambiguation
    # ----------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    apply_background(s6)
    add_header(s6, "Risk Scoring Formula & Disambiguation Matrix")

    add_card(s6, 0.8, 1.8, 11.7, 1.5, "Composite Risk Score Formula", [
        "Risk Score = 0.65 * Supervised_Prob + 0.35 * Anomaly_Score",
        "Supervised and Unsupervised components are reported separately to enable precise decision support."
    ])

    add_card(s6, 0.8, 3.6, 5.6, 3.2, "High Supervised + High Anomaly", [
        "CLASSIFICATION: High-Confidence Theft Flag",
        "ACTION: Immediate dispatch for physical inspection.",
        "SIGNATURE: Known tampering pattern with strong statistical outlier metrics."
    ])

    add_card(s6, 6.9, 3.6, 5.6, 3.2, "Low Supervised + High Anomaly", [
        "CLASSIFICATION: Novel / Zero-Day Pattern",
        "ACTION: Flagged for Human Utility Review (Not automatic theft).",
        "SIGNATURE: Unusual usage (vacant home or new tampering method)."
    ])

    s6.notes_slide.notes_text_frame.text = (
        "Presenter Script: Our composite risk score combines supervised and anomaly signals. "
        "Reporting both components separately prevents false accusations—an account with high anomaly but low supervised score "
        "is triaged as 'needs human review' rather than a definitive theft call."
    )

    # ----------------------------------------------------
    # SLIDE 7: BACKEND XAI vs FRONTEND DASHBOARD UX
    # ----------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    apply_background(s7)
    add_header(s7, "Explainability: Backend Calculation vs. Frontend UI")

    add_card(s7, 0.8, 1.8, 5.6, 5.0, "⚙️ Backend SHAP Computations", [
        "TreeExplainer Pipeline: Computes game-theoretic Shapley values per feature for every prediction.",
        "Attribution Ranking: Ranks features by absolute contribution to the risk score.",
        "Rule-Based Natural Language Translation: Maps raw SHAP values into clean operator reasons (e.g. 'Long zero streak detected').",
        "JSON Response Payload: Formats reasons array into clean REST API response."
    ], title_color=ACCENT_BLUE)

    add_card(s7, 6.9, 1.8, 5.6, 5.0, "🖥️ Frontend Operator Experience", [
        "Automated Audit Reasons Selector: Interactive dropdown allowing instant selection of flagged accounts.",
        "Visual Reason Alerts: Renders red alert callouts (🛑) highlighting exact failure modes.",
        "Transformer Zone Context: Displays transformer ID and percentage calculated risk.",
        "Grid Plotly Distribution: Renders suspect counts per transformer zone in a visual heatmap bar chart."
    ], title_color=ACCENT_GREEN)

    s7.notes_slide.notes_text_frame.text = (
        "Presenter Script: On Slide 7 we highlight the bridge between backend machine learning and frontend user experience. "
        "The backend calculates Shapley values and translates them to reasons; the frontend renders them as interactive audit logs for grid operators."
    )

    # ----------------------------------------------------
    # SLIDE 8: Datasets & Benchmark Validation
    # ----------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    apply_background(s8)
    add_header(s8, "Datasets: Synthetic Quickstart & Real Benchmarks")

    add_card(s8, 0.8, 1.8, 5.6, 5.0, "Synthetic Data Generator", [
        "Built-in generator for instant zero-dependency setup.",
        "Simulates normal load profiles, drop-offs, and bypasses.",
        "Used for automated CI/CD unit testing & quick demo verification.",
        "Reaches high ROC-AUC on clean benchmark scenarios."
    ])

    add_card(s8, 6.9, 1.8, 5.6, 5.0, "Real-World Public Benchmarks", [
        "SGCC Benchmark (State Grid Corp of China):",
        "  • 42,372 consumers, 3,615 confirmed theft cases (IEEE TII standard).",
        "  • Plug-and-play format via simple melt transformation.",
        "PRECON & Irish CER Datasets:",
        "  • High-frequency smart meter data for 30-minute interval analysis."
    ])

    s8.notes_slide.notes_text_frame.text = (
        "Presenter Script: We provide an out-of-the-box synthetic generator for quick deployment, "
        "and our codebase drops directly into standard industry benchmarks like the 42,000-consumer SGCC dataset."
    )

    # ----------------------------------------------------
    # SLIDE 9: Honest Caveats & Q&A Preparedness
    # ----------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    apply_background(s9)
    add_header(s9, "Honest Caveats & Q&A Readiness")

    add_card(s9, 0.8, 1.8, 5.6, 5.0, "Technical & Data Realities", [
        "Synthetic vs. Real Noise: Synthetic data is cleanly separable (~1.0 AUC); real-world grid data is significantly messier.",
        "Topology Reliance: Peer comparisons rely on accurate transformer mapping. Bad feeder data affects z-scores.",
        "Vacant Property Outliers: Anomaly detection flags vacant/relocated homes—hence the need for human review."
    ])

    add_card(s9, 6.9, 1.8, 5.6, 5.0, "Deployment & Policy Guidelines", [
        "Decision Support Tool: The system assists human inspectors; it is NOT an automated legal/billing penalty engine.",
        "Feedback Loop: Field inspection results feed back into retraining the supervised XGBoost model.",
        "Privacy & Governance: Operates on anonymized consumption series."
    ])

    s9.notes_slide.notes_text_frame.text = (
        "Presenter Script: In a real pitch, honesty builds trust. Real grid data is messier than synthetic demos, "
        "peer comparison requires correct feeder topology, and our tool is designed for human decision support, not automatic billing fines."
    )

    # ----------------------------------------------------
    # SLIDE 10: TECHNICAL CONTRIBUTIONS SUMMARY (BACKEND vs FRONTEND)
    # ----------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    apply_background(s10)
    add_header(s10, "Summary of Technical Contributions")

    add_card(s10, 0.8, 1.8, 5.6, 5.0, "⚙️ Backend Core Accomplishments", [
        "FastAPI REST API: Scalable microservice with POST /scan & POST /score endpoints.",
        "Dual ML Pipeline: Class-weighted XGBoost + Isolation Forest anomaly detection.",
        "15-Feature Extractor: Automated time-series dynamics, bypass streaks & peer z-scores.",
        "SHAP XAI Engine: TreeExplainer integration producing per-account risk explanations."
    ], title_color=ACCENT_BLUE)

    add_card(s10, 6.9, 1.8, 5.6, 5.0, "🖥️ Frontend UX Accomplishments", [
        "Streamlit Dashboard: Low-latency operator decision dashboard.",
        "Real-Time Parameterization: Interactive suspect sliders and dynamic scan triggers.",
        "Transformer Heatmap: Plotly visual aggregation of suspect clusters by transformer feeder.",
        "Audit Inspector: Transparent audit reason renderer restoring inspector trust."
    ], title_color=ACCENT_GREEN)

    s10.notes_slide.notes_text_frame.text = (
        "Presenter Script: To summarize our technical achievements: Our backend brings robust ML engineering, "
        "feature extraction, and API design. Our frontend delivers a slick, operator-centric Streamlit interface that turns complex AI outputs into actionable grid insights. Thank you!"
    )

    output_path = "Electricity_Theft_Detection_Deck.pptx"
    prs.save(output_path)
    print(f"Presentation updated with Backend/Frontend separation at {output_path}")

if __name__ == "__main__":
    create_deck()
